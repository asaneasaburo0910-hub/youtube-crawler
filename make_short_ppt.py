import csv
import os
import requests
import qrcode
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO

CSV_FILE = "youtube_trends.csv"
OUTPUT_FILE = f"youtube_short_{datetime.now().strftime('%Y%m%d')}.pptx"

# 縦長（9:16）サイズ
ZUNDAMON_GIF = "zundamon-an.gif"

# 縦長（9:16）サイズ
SLIDE_W = Inches(6.0)
SLIDE_H = Inches(10.67)

BG_DARK  = RGBColor(0x0F, 0x0F, 0x0F)
RED      = RGBColor(0xFF, 0x00, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xAA, 0xAA, 0xAA)
GOLD     = RGBColor(0xFF, 0xD7, 0x00)
SILVER   = RGBColor(0xC0, 0xC0, 0xC0)
BRONZE   = RGBColor(0xCD, 0x7F, 0x32)
TEAL     = RGBColor(0x4E, 0xC9, 0xB0)
RANK_COLORS = [GOLD, SILVER, BRONZE]


def make_qr(url):
    qr = qrcode.QRCode(version=1, box_size=6, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def load_csv():
    if not os.path.exists(CSV_FILE):
        print(f"❌ {CSV_FILE} が見つかりません")
        return []

    with open(CSV_FILE, newline="", encoding="utf-8") as f:
        rows = list(csv.DictReader(f))

    if not rows:
        return []

    latest_time = max(r["fetched_at"] for r in rows)
    trending = [r for r in rows if r["fetched_at"] == latest_time and r["type"] == "急上昇動画"]
    trending.sort(key=lambda x: int(x["rank"]))
    return trending[:5]


def get_thumbnail(url):
    try:
        video_id = url.split("v=")[-1].split("&")[0]
        # ショート用に高解像度サムネイルを取得
        for quality in ["hqdefault", "mqdefault"]:
            res = requests.get(f"https://img.youtube.com/vi/{video_id}/{quality}.jpg", timeout=10)
            if res.status_code == 200:
                return BytesIO(res.content)
    except Exception:
        pass
    return None


def add_bg(slide, prs):
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()


def add_text(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.CENTER, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return tf



def add_zundamon(slide):
    """右下にずんだもんGIFを挿入"""
    if not os.path.exists(ZUNDAMON_GIF):
        print("⚠️ zundamon-an.gif が見つかりません")
        return
    try:
        size = Inches(2.2)
        x = SLIDE_W - size - Inches(0.05)
        y = SLIDE_H - size - Inches(0.05)
        slide.shapes.add_picture(ZUNDAMON_GIF, x, y, size, size)
    except Exception as e:
        print(f"⚠️ ずんだもん挿入失敗: {e}")

def make_title_slide(prs, date_str):
    """タイトルスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    # 上部赤帯
    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.5), SLIDE_W, Inches(2.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    # YouTubeアイコン風▶
    add_text(slide, "▶", Inches(0), Inches(1.5), SLIDE_W, Inches(1.5),
             72, WHITE, bold=True)

    # タイトル
    add_text(slide, "YouTube\n急上昇ランキング", Inches(0.2), Inches(3.6), Inches(5.6), Inches(1.6),
             32, WHITE, bold=True)

    # TOP5
    add_text(slide, "TOP 5", Inches(0.2), Inches(5.3), SLIDE_W - Inches(0.4), Inches(0.8),
             48, GOLD, bold=True)

    # 日付
    add_text(slide, date_str, Inches(0.2), Inches(6.2), SLIDE_W - Inches(0.4), Inches(0.5),
             18, GRAY)

    # ずんだもん用スペース（下部）
    add_text(slide, "ずんだもんと一緒に見ていくのだ！", Inches(0.2), Inches(7.5), SLIDE_W - Inches(0.4), Inches(1.0),
             18, TEAL, bold=True)
    add_zundamon(slide)


def make_rank_slide(prs, item, rank):
    """ランキングスライド（縦長）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    rank_color = RANK_COLORS[rank - 1] if rank <= 3 else WHITE

    # ランク番号（大きく上部）
    add_text(slide, f"#{rank}", Inches(0), Inches(0.3), SLIDE_W, Inches(1.5),
             80, rank_color, bold=True)

    # サムネイル（中央・正方形気味）
    thumb = get_thumbnail(item["url"])
    thumb_w = Inches(5.6)
    thumb_h = Inches(3.15)  # 16:9比率
    thumb_x = (SLIDE_W - thumb_w) / 2
    if thumb:
        try:
            slide.shapes.add_picture(thumb, thumb_x, Inches(1.9), thumb_w, thumb_h)
        except Exception:
            pass

    # タイトル
    title = item["title"][:40] + ("…" if len(item["title"]) > 40 else "")
    add_text(slide, title, Inches(0.2), Inches(5.2), SLIDE_W - Inches(0.4), Inches(1.8),
             20, WHITE, bold=True)

    # チャンネル名
    add_text(slide, f"📺 {item['channel']}", Inches(0.2), Inches(7.1), SLIDE_W - Inches(0.4), Inches(0.5),
             15, GRAY)

    # 再生数・いいね
    views = int(item["views"])
    likes = int(item.get("likes", 0))
    views_str = f"{views // 10000}万" if views >= 10000 else f"{views:,}"
    likes_str = f"{likes // 10000}万" if likes >= 10000 else f"{likes:,}"
    add_text(slide, f"👁 {views_str}回　👍 {likes_str}",
             Inches(0.2), Inches(7.65), SLIDE_W - Inches(0.4), Inches(0.5),
             18, rank_color, bold=True)

    # QRコード（右下）
    try:
        qr_buf = make_qr(item["url"])
        slide.shapes.add_picture(qr_buf, Inches(0.1), Inches(8.8), Inches(1.4), Inches(1.4))
    except Exception:
        pass

    # URL
    add_text(slide, "▲ QRコードで視聴", Inches(1.6), Inches(9.1), SLIDE_W - Inches(1.8), Inches(0.4),
             12, GRAY, align=PP_ALIGN.LEFT)
    add_zundamon(slide)


def make_first_place_slide(prs, item):
    """1位特別スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    # 金色ライン上下
    for y_pos in [Inches(0), Inches(10.55)]:
        line = slide.shapes.add_shape(1, Inches(0), y_pos, SLIDE_W, Inches(0.12))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.fill.background()

    # 🏆
    add_text(slide, "🏆 NO.1 🏆", Inches(0), Inches(0.3), SLIDE_W, Inches(1.2),
             44, GOLD, bold=True)

    # サムネイル
    thumb = get_thumbnail(item["url"])
    thumb_w = Inches(5.6)
    thumb_h = Inches(3.15)
    thumb_x = (SLIDE_W - thumb_w) / 2
    if thumb:
        try:
            slide.shapes.add_picture(thumb, thumb_x, Inches(1.7), thumb_w, thumb_h)
        except Exception:
            pass

    # タイトル
    title = item["title"][:40] + ("…" if len(item["title"]) > 40 else "")
    add_text(slide, title, Inches(0.2), Inches(5.0), SLIDE_W - Inches(0.4), Inches(1.8),
             20, WHITE, bold=True)

    # チャンネル・再生数
    views = int(item["views"])
    likes = int(item.get("likes", 0))
    views_str = f"{views // 10000}万" if views >= 10000 else f"{views:,}"
    likes_str = f"{likes // 10000}万" if likes >= 10000 else f"{likes:,}"
    add_text(slide, f"📺 {item['channel']}", Inches(0.2), Inches(6.9), SLIDE_W - Inches(0.4), Inches(0.5),
             15, GRAY)
    add_text(slide, f"👁 {views_str}回　👍 {likes_str}",
             Inches(0.2), Inches(7.45), SLIDE_W - Inches(0.4), Inches(0.55),
             20, GOLD, bold=True)

    # QRコード（中央下）
    try:
        qr_buf = make_qr(item["url"])
        slide.shapes.add_picture(qr_buf, Inches(0.1), Inches(8.5), Inches(1.8), Inches(1.8))
    except Exception:
        pass

    add_text(slide, "▲ QRコードで視聴", Inches(2.0), Inches(9.1), SLIDE_W - Inches(2.1), Inches(0.4),
             12, GRAY, align=PP_ALIGN.LEFT)
    add_zundamon(slide)


def make_outro_slide(prs):
    """アウトロスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text(slide, "以上！", Inches(0), Inches(2.0), SLIDE_W, Inches(1.2),
             52, WHITE, bold=True)
    add_text(slide, "YouTube急上昇\nTOP5でした！", Inches(0.2), Inches(3.3), SLIDE_W - Inches(0.4), Inches(2.0),
             34, RED, bold=True)
    add_text(slide, "チャンネル登録・\nいいね！よろしくなのだ🌿", Inches(0.2), Inches(6.0), SLIDE_W - Inches(0.4), Inches(2.0),
             24, TEAL, bold=True)
    add_text(slide, "また明日もお楽しみに！", Inches(0.2), Inches(8.5), SLIDE_W - Inches(0.4), Inches(0.8),
             20, GRAY)
    add_zundamon(slide)


def main():
    print(f"🚀 ショートPPT生成開始: {datetime.now().strftime('%Y-%m-%d %H:%M')}")

    trending = load_csv()
    if not trending:
        print("❌ データがありません")
        return

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    latest_date = trending[0]["fetched_at"][:10]
    dt = datetime.strptime(latest_date, "%Y-%m-%d")
    date_str = dt.strftime("%Y年%m月%d日")

    # 1. タイトル
    make_title_slide(prs, date_str)

    # 2. 5位→2位（降順）
    for item in sorted(trending[1:], key=lambda x: int(x["rank"]), reverse=True):
        make_rank_slide(prs, item, int(item["rank"]))

    # 3. 1位（特別）
    make_first_place_slide(prs, trending[0])

    # 4. アウトロ
    make_outro_slide(prs)

    prs.save(OUTPUT_FILE)
    print(f"✅ {OUTPUT_FILE} を生成しました（{len(prs.slides)}スライド）")
    print(f"   構成：タイトル→5位→4位→3位→2位→1位→アウトロ")


if __name__ == "__main__":
    main()
