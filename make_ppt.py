import csv
import os
import requests
import qrcode
from datetime import datetime, timezone, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO

CSV_FILE = "youtube_trends.csv"
OUTPUT_FILE = f"youtube_ranking_{datetime.now().strftime('%Y%m%d')}.pptx"

ZUNDAMON_GIF = "zundamon-an.gif"

BG_DARK  = RGBColor(0x0F, 0x0F, 0x0F)
BG_CARD  = RGBColor(0x1E, 0x1E, 0x1E)
RED      = RGBColor(0xFF, 0x00, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xAA, 0xAA, 0xAA)
GOLD     = RGBColor(0xFF, 0xD7, 0x00)
SILVER   = RGBColor(0xC0, 0xC0, 0xC0)
BRONZE   = RGBColor(0xCD, 0x7F, 0x32)
TEAL     = RGBColor(0x4E, 0xC9, 0xB0)
RANK_COLORS = [GOLD, SILVER, BRONZE]


def make_qr(url):
    """URLからQRコードのBytesIOを生成"""
    qr = qrcode.QRCode(version=1, box_size=6, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def load_csv():
    if not os.path.exists(CSV_FILE):
        print(f"❌ {CSV_FILE} が見つかりません")
        return [], {}

    with open(CSV_FILE, newline="", encoding="utf-8") as f:
        rows = list(csv.DictReader(f))

    if not rows:
        return [], {}

    latest_time = max(r["fetched_at"] for r in rows)

    prev_rows = [r for r in rows if r["fetched_at"] != latest_time and r["type"] == "急上昇動画"]
    prev_views = {}
    if prev_rows:
        prev_time = max(r["fetched_at"] for r in prev_rows)
        for r in prev_rows:
            if r["fetched_at"] == prev_time:
                video_id = r["url"].split("v=")[-1].split("&")[0]
                prev_views[video_id] = int(r["views"])

    trending = []
    keyword_items = {}

    for row in rows:
        if row["fetched_at"] != latest_time:
            continue
        if row["type"] == "急上昇動画":
            video_id = row["url"].split("v=")[-1].split("&")[0]
            prev = prev_views.get(video_id)
            views = int(row["views"])
            row["growth_rate"] = round((views - prev) / prev * 100, 1) if prev and prev > 0 else None
            trending.append(row)
        elif row["type"] == "キーワード検索":
            label = row["label"]
            if label not in keyword_items:
                keyword_items[label] = []
            keyword_items[label].append(row)

    trending.sort(key=lambda x: int(x["rank"]))
    print("📋 ランク順:", [int(x["rank"]) for x in trending[:10]])
    return trending, keyword_items


def hours_since_published(published_at):
    try:
        dt = datetime.strptime(published_at, "%Y-%m-%d %H:%M")
        return int((datetime.now() - dt).total_seconds() / 3600)
    except Exception:
        return None


def generate_comment(item):
    comments = []
    views = int(item["views"])
    likes = int(item.get("likes", 0))
    comment_count = int(item.get("comments", 0))
    growth = item.get("growth_rate")
    published_at = item.get("published_at", "")
    hours = hours_since_published(published_at)

    if hours is not None:
        if hours <= 24:
            comments.append(f"公開{hours}時間で急上昇入り！")
        elif hours <= 72:
            comments.append(f"公開{hours//24}日で急上昇")

    if views >= 5_000_000:
        comments.append("500万再生超の超バズ動画")
    elif views >= 1_000_000:
        comments.append("100万再生超の大ヒット")
    elif views >= 500_000:
        comments.append("50万再生の注目コンテンツ")
    elif views >= 100_000:
        comments.append("10万再生突破の話題作")
    else:
        comments.append("急上昇中の新興コンテンツ")

    if growth is not None:
        if growth >= 200:
            comments.append(f"前日比+{growth}%の爆発的な伸び")
        elif growth >= 50:
            comments.append(f"前日比+{growth}%と急拡散中")
        elif growth >= 10:
            comments.append(f"前日比+{growth}%で着実に成長")

    if views > 0:
        like_rate = likes / views * 100
        if like_rate >= 5:
            comments.append(f"いいね率{like_rate:.1f}%と視聴者の支持が非常に高い")
        elif like_rate >= 2:
            comments.append(f"いいね率{like_rate:.1f}%と好評価")

    if comment_count >= 10000:
        comments.append("コメント1万超の大反響")
    elif comment_count >= 1000:
        comments.append("コメント1000件超と議論が活発")

    return "　".join(comments[:3])


def get_thumbnail(url):
    try:
        video_id = url.split("v=")[-1].split("&")[0]
        res = requests.get(f"https://img.youtube.com/vi/{video_id}/mqdefault.jpg", timeout=10)
        if res.status_code == 200:
            return BytesIO(res.content)
    except Exception:
        pass
    return None


def add_bg(slide, prs, color=None):
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color or BG_DARK
    bg.line.fill.background()


def add_text(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, name="Arial", wrap=False):
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    return tf


def add_text_with_link(slide, text, url, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, name="Arial", wrap=False):
    """ハイパーリンク付きテキスト"""
    from pptx.oxml.ns import qn
    from lxml import etree
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    run.font.underline = True
    rPr = run._r.get_or_add_rPr()
    hlinkClick = etree.SubElement(rPr, qn('a:hlinkClick'))
    rId = slide.part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
    hlinkClick.set(qn('r:id'), rId)
    return tf



def add_zundamon(slide):
    """右下にずんだもんGIFを挿入"""
    if not os.path.exists(ZUNDAMON_GIF):
        print("⚠️ zundamon-an.gif が見つかりません")
        return
    try:
        size = Inches(1.5)
        from pptx.util import Inches as I
        x = Inches(10) - size - Inches(0.05)
        y = Inches(5.625) - size - Inches(0.05)
        slide.shapes.add_picture(ZUNDAMON_GIF, x, y, size, size)
    except Exception as e:
        print(f"⚠️ ずんだもん挿入失敗: {e}")

def make_title_slide(prs, date_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.1), prs.slide_width, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    add_text(slide, "▶  YouTube トレンド ランキング", Inches(0.5), Inches(2.18), Inches(9), Inches(1.0),
             34, WHITE, bold=True, align=PP_ALIGN.CENTER, name="Arial Black")
    add_text(slide, f"{date_str}  ·  急上昇 TOP 10 発表", Inches(0.5), Inches(3.6), Inches(9), Inches(0.6),
             15, GRAY, align=PP_ALIGN.CENTER)
    add_zundamon(slide)


def make_rank_slide(prs, item, rank):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    rank_color = RANK_COLORS[rank - 1] if rank <= 3 else WHITE

    # ランク番号
    add_text(slide, f"#{rank}", Inches(0.35), Inches(0.1), Inches(2), Inches(0.85),
             48, rank_color, bold=True, name="Arial Black")

    # サムネイル
    thumb = get_thumbnail(item["url"])
    if thumb:
        try:
            slide.shapes.add_picture(thumb, Inches(0.35), Inches(1.05), Inches(4.0), Inches(2.25))
        except Exception:
            pass

    # QRコード（サムネイル右下）
    try:
        qr_buf = make_qr(item["url"])
        slide.shapes.add_picture(qr_buf, Inches(0.35), Inches(3.38), Inches(1.1), Inches(1.1))
    except Exception:
        pass

    # URL テキスト（QRの右）
    add_text(slide, item["url"], Inches(1.55), Inches(3.55), Inches(2.85), Inches(0.4),
             7, GRAY, wrap=False)

    # タイトル（リンク付き）
    add_text_with_link(slide, item["title"], item["url"], Inches(4.55), Inches(0.85), Inches(5.1), Inches(1.6),
                       14, WHITE, bold=True, wrap=True)

    # チャンネル名
    add_text(slide, f"📺  {item['channel']}", Inches(4.55), Inches(2.55), Inches(5.1), Inches(0.32),
             11, GRAY)

    # 再生数・いいね・コメント
    views = int(item["views"])
    likes = int(item.get("likes", 0))
    comments = int(item.get("comments", 0))
    growth = item.get("growth_rate")
    growth_str = f"  📈+{growth}%" if growth and growth > 0 else ""
    add_text(slide, f"👁 {views:,}  👍 {likes:,}  💬 {comments:,}{growth_str}",
             Inches(4.55), Inches(2.92), Inches(5.1), Inches(0.35), 11, rank_color, bold=True)

    # 動画の長さ・公開日時
    duration = item.get("duration", "不明")
    published_at = item.get("published_at", "")
    hours = hours_since_published(published_at)
    hours_str = f"公開{hours}時間前" if hours is not None and hours <= 168 else f"公開 {published_at[:10]}"
    add_text(slide, f"⏱ {duration}  ·  📅 {hours_str}",
             Inches(4.55), Inches(3.32), Inches(5.1), Inches(0.32), 11, GRAY)

    # 分析コメント
    comment = generate_comment(item)
    add_text(slide, f"💡 {comment}", Inches(0.35), Inches(4.55), Inches(9.3), Inches(0.6),
             11, TEAL, wrap=True)
    add_zundamon(slide)


def make_first_place_slide(prs, item):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, RGBColor(0x08, 0x08, 0x08))

    for y_pos in [Inches(0), Inches(5.5)]:
        line = slide.shapes.add_shape(1, Inches(0), y_pos, prs.slide_width, Inches(0.12))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.fill.background()

    add_text(slide, "🏆  NO. 1  🏆", Inches(0), Inches(0.12), Inches(10), Inches(0.95),
             36, GOLD, bold=True, align=PP_ALIGN.CENTER, name="Arial Black")

    # サムネイル（中央）
    thumb = get_thumbnail(item["url"])
    if thumb:
        try:
            thumb_w, thumb_h = Inches(5.0), Inches(2.81)
            slide.shapes.add_picture(thumb, (Inches(10) - thumb_w) / 2, Inches(1.12), thumb_w, thumb_h)
        except Exception:
            pass

    # QRコード（右下）
    try:
        qr_buf = make_qr(item["url"])
        slide.shapes.add_picture(qr_buf, Inches(0.1), Inches(3.98), Inches(1.1), Inches(1.1))
    except Exception:
        pass

    # タイトル（リンク付き）
    title = item["title"][:55] + ("…" if len(item["title"]) > 55 else "")
    add_text_with_link(slide, title, item["url"], Inches(0.3), Inches(4.0), Inches(8.1), Inches(0.5),
                       14, WHITE, bold=True, align=PP_ALIGN.CENTER, wrap=True)

    # 再生数・いいね・コメント・動画長・公開日時
    views = int(item["views"])
    likes = int(item.get("likes", 0))
    comments_count = int(item.get("comments", 0))
    duration = item.get("duration", "")
    published_at = item.get("published_at", "")
    hours = hours_since_published(published_at)
    hours_str = f"公開{hours}時間前" if hours is not None and hours <= 168 else published_at[:10]
    add_text(slide, f"👁 {views:,}  👍 {likes:,}  💬 {comments_count:,}  ⏱ {duration}  📅 {hours_str}",
             Inches(0.3), Inches(4.55), Inches(8.1), Inches(0.32), 11, GOLD, align=PP_ALIGN.CENTER)

    # 分析コメント
    comment = generate_comment(item)
    add_text(slide, f"💡 {comment}", Inches(0.3), Inches(4.92), Inches(8.1), Inches(0.32),
             10, TEAL, align=PP_ALIGN.CENTER, wrap=True)

    # URL テキスト
    add_text(slide, item["url"], Inches(0.3), Inches(5.28), Inches(8.1), Inches(0.25),
             7, GRAY, align=PP_ALIGN.CENTER)
    add_zundamon(slide)


def make_keyword_slide(prs, label, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    emoji = {"エンタメ": "🎬", "ゲーム": "🎮", "アニメ": "✨", "映画": "🍿"}.get(label, "📌")
    add_text(slide, f"{emoji} {label} 人気動画 TOP {min(5, len(items))}",
             Inches(0.4), Inches(0.15), Inches(9), Inches(0.55), 22, WHITE, bold=True, name="Arial Black")

    card_h = Inches(0.88)
    for i, item in enumerate(items[:5]):
        rank = int(item["rank"])
        y = Inches(0.82) + i * (card_h + Inches(0.08))

        card = slide.shapes.add_shape(1, Inches(0.3), y, Inches(9.4), card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_CARD
        card.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        card.line.width = Pt(0.5)

        thumb = get_thumbnail(item["url"])
        if thumb:
            try:
                slide.shapes.add_picture(thumb, Inches(0.35), y, Inches(1.9), Inches(1.07))
            except Exception:
                pass

        rank_color = RANK_COLORS[rank - 1] if rank <= 3 else GRAY
        add_text(slide, f"#{rank}", Inches(2.4), y + Inches(0.04), Inches(0.6), Inches(0.4),
                 15, rank_color, bold=True, name="Arial Black")

        title = item["title"][:42] + ("…" if len(item["title"]) > 42 else "")
        add_text_with_link(slide, title, item["url"], Inches(2.95), y + Inches(0.04), Inches(5.7), Inches(0.38),
                           12, WHITE, bold=True, wrap=True)

        # QRコード（右端）
        try:
            qr_buf = make_qr(item["url"])
            slide.shapes.add_picture(qr_buf, Inches(8.8), y + Inches(0.06), Inches(0.75), Inches(0.75))
        except Exception:
            pass

        views = int(item["views"])
        likes = int(item.get("likes", 0))
        duration = item.get("duration", "")
        published_at = item.get("published_at", "")[:10]
        add_text(slide, f"📺 {item['channel']}  👁 {views:,}  👍 {likes:,}  ⏱ {duration}  📅 {published_at}",
                 Inches(2.95), y + Inches(0.50), Inches(5.7), Inches(0.32), 10, GRAY)
        # URL テキスト
        add_text(slide, item["url"], Inches(2.95), y + Inches(0.70), Inches(5.7), Inches(0.22),
                 7, RGBColor(0x66, 0x66, 0x66))


def main():
    print(f"🚀 PPT生成開始: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    trending, keyword_items = load_csv()
    if not trending and not keyword_items:
        print("❌ データがありません")
        return

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    make_title_slide(prs, datetime.now().strftime("%Y年%m月%d日"))

    top10 = trending[:10]
    for item in sorted(top10[1:], key=lambda x: int(x["rank"]), reverse=True):
        make_rank_slide(prs, item, int(item["rank"]))

    if top10:
        make_first_place_slide(prs, top10[0])

    for label, items in keyword_items.items():
        make_keyword_slide(prs, label, items)

    prs.save(OUTPUT_FILE)
    print(f"✅ {OUTPUT_FILE} を生成しました（{len(prs.slides)}スライド）")


if __name__ == "__main__":
    main()
